"""Render a .pptx from a JSON spec. JSON defines layout; images live in image_dir."""
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

def hex_color(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def set_text(tf, content, font="Calibri", size=14, bold=False, italic=False,
             color="2B2B2B", letter_spacing=0, line_spacing=None, align=None):
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = ""
    if align == "center": p.alignment = PP_ALIGN.CENTER
    elif align == "right": p.alignment = PP_ALIGN.RIGHT
    if line_spacing: p.line_spacing = line_spacing
    run = p.add_run()
    run.text = content
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_color(color)
    if letter_spacing:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(letter_spacing * 100))

def add_text(slide, el):
    box = slide.shapes.add_textbox(
        Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"])
    )
    set_text(
        box.text_frame, el["content"],
        font=el.get("font","Calibri"), size=el.get("size",14),
        bold=el.get("bold",False), italic=el.get("italic",False),
        color=el.get("color","2B2B2B"),
        letter_spacing=el.get("letter_spacing",0),
        line_spacing=el.get("line_spacing"),
        align=el.get("align"),
    )

def add_rect(slide, el):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"])
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_color(el["fill"])
    shape.line.fill.background()
    shape.shadow.inherit = False

def add_image(slide, el, image_dir):
    path = os.path.join(image_dir, el["file"])
    pic = slide.shapes.add_picture(
        path, Inches(el["x"]), Inches(el["y"]),
        Inches(el["w"]), Inches(el["h"])
    )
    if el.get("rounded"):
        sp = pic._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is not None:
            prstGeom = spPr.find(qn("a:prstGeom"))
            if prstGeom is not None:
                prstGeom.set("prst", "roundRect")

def add_flow_node(slide, el):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if el.get("shape")=="rounded_rect" else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(
        shape_type,
        Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"])
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_color(el.get("fill","FFFFFF"))
    shp.line.color.rgb = hex_color(el.get("border","666666"))
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.18); tf.margin_bottom = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = ""
    r = p.add_run(); r.text = el["title"]
    r.font.name = "Georgia"; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = hex_color(el.get("title_color","2B2B2B"))
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    r2 = p2.add_run(); r2.text = el["body"]
    r2.font.name = "Calibri"; r2.font.size = Pt(11)
    r2.font.color.rgb = hex_color(el.get("body_color","2B2B2B"))

def add_flow_arrow(slide, el):
    x1, y1 = el["from"]; x2, y2 = el["to"]
    w = abs(x2 - x1); h = 0.18
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(min(x1,x2)), Inches(y1 - h/2), Inches(w), Inches(h)
    )
    shp.fill.solid(); shp.fill.fore_color.rgb = hex_color(el.get("color","999999"))
    shp.line.fill.background()
    shp.shadow.inherit = False

def add_stat(slide, el):
    box = slide.shapes.add_textbox(
        Inches(el["x"]), Inches(el["y"]), Inches(el["w"]), Inches(el["h"])
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = ""
    r = p.add_run(); r.text = el["number"]
    r.font.name = "Georgia"; r.font.size = Pt(54); r.font.bold = True
    r.font.color.rgb = hex_color("6D2E46")
    if el.get("unit"):
        r2 = p.add_run(); r2.text = "  " + el["unit"]
        r2.font.name = "Calibri"; r2.font.size = Pt(16)
        r2.font.color.rgb = hex_color("A26769")
    p2 = tf.add_paragraph()
    p2.space_before = Pt(2)
    r3 = p2.add_run(); r3.text = el["label"]
    r3.font.name = "Calibri"; r3.font.size = Pt(12)
    r3.font.color.rgb = hex_color("7A6A6A")

DISPATCH = {
    "rect": lambda s,e,d: add_rect(s,e),
    "text": lambda s,e,d: add_text(s,e),
    "image": lambda s,e,d: add_image(s,e,d),
    "flow_node": lambda s,e,d: add_flow_node(s,e),
    "flow_arrow": lambda s,e,d: add_flow_arrow(s,e),
    "stat": lambda s,e,d: add_stat(s,e),
}

def render(spec_path, out_path):
    base = os.path.dirname(os.path.abspath(spec_path))
    spec = json.load(open(spec_path, encoding="utf-8"))
    image_dir = os.path.join(base, spec["meta"].get("image_dir","images"))
    prs = Presentation()
    sw = spec["meta"]["slide_size"]["width"]
    sh = spec["meta"]["slide_size"]["height"]
    prs.slide_width = Inches(sw); prs.slide_height = Inches(sh)
    blank = prs.slide_layouts[6]
    bg_hex = spec["theme"].get("background","FFFFFF")
    for slide_spec in spec["slides"]:
        slide = prs.slides.add_slide(blank)
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = hex_color(bg_hex)
        bg.line.fill.background()
        bg.shadow.inherit = False
        for el in slide_spec["elements"]:
            fn = DISPATCH.get(el["type"])
            if not fn:
                print("unknown element:", el["type"]); continue
            fn(slide, el, image_dir)
    prs.save(out_path)
    print("wrote", out_path)

if __name__ == "__main__":
    spec = sys.argv[1] if len(sys.argv)>1 else "spec.json"
    out  = sys.argv[2] if len(sys.argv)>2 else "output.pptx"
    render(spec, out)
